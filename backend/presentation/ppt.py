from pptx import Presentation


def create_ppt(presentation_json, output_file="output.pptx"):
    prs = Presentation()

    slides = presentation_json.get("slides", [])

    for slide_data in slides:
        slide_layout = prs.slide_layouts[1]  # title + content
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = slide_data.get("heading", "")

        points = slide_data.get("points", [])
        content.text = "\n".join(points)

    prs.save(output_file)

    return output_file